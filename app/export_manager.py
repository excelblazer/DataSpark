import pandas as pd
import json
import os
from typing import Dict, List, Union
from fastapi import HTTPException
import matplotlib.pyplot as plt
from io import BytesIO
import base64
from docx import Document
from pptx import Presentation
from pptx.util import Inches

class ExportManager:
    def __init__(self):
        self.supported_formats = ["csv", "excel", "json", "pdf", "ppt", "word", "powerbi", "html"]
    
    def export(self, data: Union[Dict, pd.DataFrame], format_type: str, filename: str = None) -> Dict:
        """
        Export data to the requested format
        
        Args:
            data: Dictionary or DataFrame containing the data to export
            format_type: Export format (csv, excel, json, pdf, ppt, word, powerbi)
            filename: Optional filename for the exported file
            
        Returns:
            Dictionary containing the export result (file path or content)
        """
        if format_type not in self.supported_formats:
            raise HTTPException(status_code=400, detail=f"Format {format_type} not supported")
        
        # Ensure data is a DataFrame
        if isinstance(data, dict):
            try:
                # If it's a dict with visualization data
                if "figure" in data:
                    fig_data = data["figure"]
                    # Handle visualization export separately
                    return self._export_visualization(fig_data, format_type, filename)
                # If it's a dict with tabular data
                elif "data" in data:
                    df = pd.DataFrame(data["data"])
                else:
                    df = pd.DataFrame(data)
            except Exception as e:
                raise HTTPException(status_code=400, detail=f"Error converting data to DataFrame: {str(e)}")
        else:
            df = data
            
        # Export based on format
        if format_type == "csv":
            return self._export_csv(df, filename)
        elif format_type == "excel":
            return self._export_excel(df, filename)
        elif format_type == "json":
            return self._export_json(df, filename)
        elif format_type == "pdf":
            return self._export_pdf(df, filename)
        elif format_type == "ppt":
            return self._export_ppt(df, data.get("title", "Data Export"), filename)
        elif format_type == "word":
            return self._export_word(df, data.get("title", "Data Export"), filename)
        elif format_type == "powerbi":
            return self._export_powerbi(df, filename)
        elif format_type == "html":
            return self._export_html(df, filename)
            
    def _export_csv(self, df: pd.DataFrame, filename: str = None) -> Dict:
        """Export data to CSV"""
        if filename is None:
            filename = f"export_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.csv"
        
        filepath = f"./exports/{filename}"
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        
        df.to_csv(filepath, index=False)
        return {"format": "csv", "filepath": filepath, "filename": filename}
    
    def _export_excel(self, df: pd.DataFrame, filename: str = None) -> Dict:
        """Export data to Excel"""
        if filename is None:
            filename = f"export_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        
        filepath = f"./exports/{filename}"
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        
        df.to_excel(filepath, index=False)
        return {"format": "excel", "filepath": filepath, "filename": filename}
    
    def _export_json(self, df: pd.DataFrame, filename: str = None) -> Dict:
        """Export data to JSON"""
        if filename is None:
            filename = f"export_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.json"
        
        filepath = f"./exports/{filename}"
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        
        df.to_json(filepath, orient="records")
        return {"format": "json", "filepath": filepath, "filename": filename}
    
    def _export_pdf(self, df: pd.DataFrame, filename: str = None) -> Dict:
        """Export data to PDF"""
        if filename is None:
            filename = f"export_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        
        filepath = f"./exports/{filename}"
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        
        # Create PDF using matplotlib and save to file
        fig, ax = plt.subplots(figsize=(12, 4))
        ax.axis('tight')
        ax.axis('off')
        table = ax.table(cellText=df.values, colLabels=df.columns, loc='center')
        
        plt.savefig(filepath, format='pdf', bbox_inches='tight')
        plt.close()
        
        return {"format": "pdf", "filepath": filepath, "filename": filename}
    
    def _export_word(self, df: pd.DataFrame, title: str, filename: str = None) -> Dict:
        """Export data to Word document"""
        if filename is None:
            filename = f"export_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.docx"
        
        filepath = f"./exports/{filename}"
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        
        doc = Document()
        doc.add_heading(title, level=1)
        
        # Add timestamp
        doc.add_paragraph(f"Generated on {pd.Timestamp.now().strftime('%Y-%m-%d %H:%M:%S')}")
        
        # Add table
        table = doc.add_table(rows=1, cols=len(df.columns))
        table.style = 'Table Grid'
        
        # Add header row
        header_cells = table.rows[0].cells
        for i, column in enumerate(df.columns):
            header_cells[i].text = str(column)
        
        # Add data rows
        for _, row in df.iterrows():
            row_cells = table.add_row().cells
            for i, value in enumerate(row):
                row_cells[i].text = str(value)
        
        doc.save(filepath)
        return {"format": "word", "filepath": filepath, "filename": filename}
    
    def _export_ppt(self, df: pd.DataFrame, title: str, filename: str = None) -> Dict:
        """Export data to PowerPoint presentation"""
        if filename is None:
            filename = f"export_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        
        filepath = f"./exports/{filename}"
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title and content layout
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Create a table for the data
        rows, cols = len(df) + 1, len(df.columns)  # +1 for header
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Add header row
        for i, column in enumerate(df.columns):
            table.cell(0, i).text = str(column)
        
        # Add data rows
        for i, (_, row) in enumerate(df.iterrows(), start=1):
            for j, value in enumerate(row):
                table.cell(i, j).text = str(value)
        
        prs.save(filepath)
        return {"format": "ppt", "filepath": filepath, "filename": filename}
    
    def _export_powerbi(self, df: pd.DataFrame, filename: str = None) -> Dict:
        """
        Export data in a format compatible with Power BI
        For this example, we'll create a PBIX compatible Excel file
        """
        if filename is None:
            filename = f"export_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}_powerbi.xlsx"
        
        filepath = f"./exports/{filename}"
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        
        # Create a formatted Excel file that works well with Power BI
        writer = pd.ExcelWriter(filepath, engine='xlsxwriter')
        df.to_excel(writer, sheet_name='Data', index=False)
        
        # Format the Excel sheet
        workbook = writer.book
        worksheet = writer.sheets['Data']
        
        # Add format for headers
        header_format = workbook.add_format({
            'bold': True,
            'text_wrap': True,
            'valign': 'top',
            'bg_color': '#D7E4BC',
            'border': 1
        })
        
        # Add headers with the defined format
        for col_num, value in enumerate(df.columns.values):
            worksheet.write(0, col_num, value, header_format)
        
        writer.close()
        
        return {
            "format": "powerbi",
            "filepath": filepath,
            "filename": filename,
            "message": "Import this Excel file into Power BI Desktop"
        }
    
    def _export_html(self, df: pd.DataFrame, filename: str = None) -> Dict:
        """Export data to HTML"""
        if filename is None:
            filename = f"export_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.html"
        
        filepath = f"./exports/{filename}"
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        
        html_content = df.to_html(classes='table table-striped', index=False)
        
        with open(filepath, 'w') as f:
            f.write(f"""
            <!DOCTYPE html>
            <html>
            <head>
                <title>Data Export</title>
                <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.0.2/dist/css/bootstrap.min.css" rel="stylesheet">
                <style>
                    body {{ padding: 20px; }}
                </style>
            </head>
            <body>
                <div class="container">
                    <h2>Data Export</h2>
                    <p>Generated on {pd.Timestamp.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
                    {html_content}
                </div>
            </body>
            </html>
            """)
        
        return {"format": "html", "filepath": filepath, "filename": filename}
    
    def _export_visualization(self, fig_data, format_type: str, filename: str = None) -> Dict:
        """Export visualization to specified format"""
        # This would handle the export of Plotly or other visualization data
        # Implementation depends on the visualization library used
        
        if format_type in ["pdf", "ppt", "word"]:
            # For formats that support images, convert the visualization to an image first
            # The implementation will depend on whether you're using Plotly, Matplotlib, etc.
            # For this example, let's assume fig_data is a Plotly figure JSON
            
            # Create a placeholder response for now
            return {
                "format": format_type,
                "message": f"Visualization export to {format_type} is not fully implemented"
            }
        else:
            # For formats like JSON, just return the figure data
            if filename is None:
                filename = f"vis_export_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.json"
            
            filepath = f"./exports/{filename}"
            os.makedirs(os.path.dirname(filepath), exist_ok=True)
            
            with open(filepath, 'w') as f:
                json.dump(fig_data, f)
            
            return {"format": "json", "filepath": filepath, "filename": filename}
